# -*- coding: utf-8 -*-
# @Author: Anderson
# @Date:   2019-07-03 15:36:58
# @Last Modified by:   Anderson
# @Last Modified time: 2021-04-15 02:01:45
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
import glob
import os


def set_text_frame_font(text_frame):
	for paragraph in text_frame.paragraphs:
		for run in paragraph.runs:
			if run.font.name in fonts_to_be_replaced:
				run.font.name = fonts_to_be_replaced[run.font.name]
				if run.font._rPr.find(qn('a:ea')) is not None:
					run.font._rPr.find(qn('a:ea')).set('typeface', run.font.name)
				else:
					element = run.font._rPr.makeelement(qn('a:ea'))
					element.set('typeface', run.font.name)
					run.font._rPr.append(element)
			elif run.font.name is None:
				run.font.name = '思源黑体'


def check_shape(shape):
	if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
		for shape_in_group in shape.shapes:
			check_shape(shape_in_group)
	elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
		for cell in shape.table.iter_cells():
			text_frame = cell.text_frame
			set_text_frame_font(text_frame)
	else:
		if shape.has_text_frame:
			text_frame = shape.text_frame
			set_text_frame_font(text_frame)


fonts_to_be_replaced = {
	'微软雅黑': '思源黑体 Normal',
	'Microsoft YaHei': '思源黑体 Normal',
	'等线': '思源黑体 Normal',
	'Open Sans': '思源黑体 Normal'
}

for file in glob.glob('input/*.pptx'):
	print(f'Processing file: {file}')
	prs = Presentation(file)
	for index, slide in enumerate(prs.slides):
		for shape in slide.shapes:
			check_shape(shape)

	prs.save(f'output/{os.path.basename(file)}')
